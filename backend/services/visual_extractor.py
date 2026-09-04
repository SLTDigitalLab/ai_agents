import os
import json
import base64
import logging
from pathlib import Path
from typing import Tuple, List
import fitz  # PyMuPDF
from langchain_core.messages import HumanMessage
from langchain_openai import ChatOpenAI
from langchain_core.documents import Document
from core.config import evidence_storage_dir, settings

log = logging.getLogger(__name__)

STORAGE_DIR = evidence_storage_dir()
STORAGE_DIR.mkdir(parents=True, exist_ok=True)
AUDIT_DIR = STORAGE_DIR / "audits"
AUDIT_DIR.mkdir(parents=True, exist_ok=True)

VLM_MODEL = "gpt-4o"
VLM_TEMPERATURE = 0

VISUAL_DETECTION_PROMPT = (
    "Analyze this document page/slide/chart image carefully.\n"
    "If it contains a flowchart, diagram, process flow, organizational chart, or data chart:\n"
    "1. Identify the title/heading of the visual.\n"
    "2. Transcribe all text within the diagram nodes and labels.\n"
    "3. Provide a step-by-step narrative of the process flow, decision logic, relationships, and outcomes.\n"
    "If the image ONLY contains standard text paragraphs or basic formatting without meaningful visuals, reply ONLY with 'NO_VISUAL_CONTENT'."
)


def _write_audit_record(
    doc_id: str,
    source_file: str,
    page_number: int,
    status: str,
    reason: str,
    **details,
) -> None:
    """Append a durable record for every PDF page assessed for Visual RAG."""
    record = {
        "doc_id": doc_id,
        "source_file": source_file,
        "page_number": page_number,
        "status": status,
        "reason": reason,
        **details,
    }
    audit_path = AUDIT_DIR / f"{doc_id}.jsonl"
    with audit_path.open("a", encoding="utf-8") as audit_file:
        audit_file.write(json.dumps(record, ensure_ascii=True) + "\n")


def encode_image(image_path: str) -> str:
    """Helper to convert local image files into base64 string."""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def analyze_image_with_gpt4o(image_path: str, prompt: str = VISUAL_DETECTION_PROMPT) -> str:
    """Send an image to GPT-4o for visual analysis."""
    try:
        base64_img = encode_image(image_path)
        llm = ChatOpenAI(model=VLM_MODEL, temperature=VLM_TEMPERATURE)
        
        message = HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{base64_img}"}
                }
            ]
        )
        
        response = llm.invoke([message])
        return response.content.strip()
    except Exception as e:
        log.error(f"GPT-4o analysis failed for {image_path}: {e}")
        return "NO_VISUAL_CONTENT"


def save_visual_record(
    doc_id: str,
    source_file: str,
    page_number: int,
    image_path: Path,
    visual_description: str
) -> dict:
    """Save visual record as JSON and return the record dict."""
    base64_img = encode_image(str(image_path))
    
    record = {
        "doc_id": doc_id,
        "source_file": source_file,
        "page_number": page_number,
        "image_path": str(image_path),
        "base64_image": base64_img,
        "visual_description": visual_description,
        "type": "visual_description"
    }
    
    # Save JSON metadata alongside visual assets
    json_path = STORAGE_DIR / f"{doc_id}_page_{page_number}.json"
    with open(json_path, "w") as f:
        json.dump(record, f, indent=2)
    
    return record


def create_vector_document(
    doc_id: str,
    source_file: str,
    page_number: int,
    image_path: str,
    visual_description: str,
    base64_img: str
) -> Document:
    """Create a LangChain Document from visual evidence."""
    return Document(
        page_content=f"[VISUAL ANALYSIS - Page {page_number}]\n{visual_description}",
        metadata={
            "doc_id": doc_id,
            "source": source_file,
            "page": page_number,
            "page_number": page_number,
            "type": "visual_description",
            "has_image": True,
            "image_path": image_path,
            "image_base64": base64_img,
            "source_type": "visual_evidence"
        }
    )


def process_pdf_visuals(pdf_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Scans PDF pages for visual content, extracts page renders, 
    generates visual descriptions using GPT-4o, and formats output for Vector DB indexing.
    """
    visual_records = []
    vector_documents = []
    
    try:
        doc = fitz.open(pdf_path)
        source_file = os.path.basename(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)
            drawings = page.get_drawings()

            # Check for raster images or visual vector components
            has_visual_signal = bool(images) or len(drawings) > 10
            if not has_visual_signal and not settings.VISUAL_AUDIT_ALL_PDF_PAGES:
                _write_audit_record(
                    doc_id,
                    source_file,
                    page_num + 1,
                    "skipped",
                    "no_visual_signal",
                    image_count=len(images),
                    drawing_count=len(drawings),
                )
                continue
                
            # 1. Render page to image (150 DPI balances quality and VLM speed)
            pix = page.get_pixmap(dpi=150)
            image_filename = f"{doc_id}_page_{page_num + 1}.png"
            image_path = STORAGE_DIR / image_filename
            pix.save(str(image_path))

            # 2. Analyze with GPT-4o
            visual_description = analyze_image_with_gpt4o(str(image_path))

            # Skip saving if no meaningful visual was detected
            if "NO_VISUAL_CONTENT" in visual_description:
                _write_audit_record(
                    doc_id,
                    source_file,
                    page_num + 1,
                    "skipped",
                    "vlm_reported_no_visual_content",
                    image_count=len(images),
                    drawing_count=len(drawings),
                    image_path=str(image_path),
                )
                if image_path.exists():
                    os.remove(image_path)
                continue

            # 3. Save record and create vector document
            base64_img = encode_image(str(image_path))
            record = save_visual_record(doc_id, source_file, page_num + 1, image_path, visual_description)
            vector_doc = create_vector_document(doc_id, source_file, page_num + 1, str(image_path), visual_description, base64_img)
            
            visual_records.append(record)
            vector_documents.append(vector_doc)
            _write_audit_record(
                doc_id,
                source_file,
                page_num + 1,
                "indexed",
                "visual_description_created",
                image_count=len(images),
                drawing_count=len(drawings),
                image_path=str(image_path),
            )

        doc.close()
        log.info(f"PDF visual extraction: extracted {len(vector_documents)} visuals from {pdf_path}")
    except Exception as e:
        log.error(f"PDF visual extraction failed for {pdf_path}: {e}")
    
    return visual_records, vector_documents


def process_docx_visuals(docx_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Extract images and shapes from DOCX files and analyze them.
    """
    visual_records = []
    vector_documents = []
    
    try:
        from docx import Document as DocxDocument
        from docx.oxml import parse_xml
        
        doc = DocxDocument(docx_path)
        source_file = os.path.basename(docx_path)
        page_counter = 0
        
        # Extract images from document relationships
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_data = image_part.blob
                    
                    page_counter += 1
                    image_filename = f"{doc_id}_page_{page_counter}.png"
                    image_path = STORAGE_DIR / image_filename
                    
                    with open(image_path, "wb") as f:
                        f.write(image_data)
                    
                    # Analyze with GPT-4o
                    visual_description = analyze_image_with_gpt4o(str(image_path))
                    
                    if "NO_VISUAL_CONTENT" in visual_description:
                        if image_path.exists():
                            os.remove(image_path)
                        continue
                    
                    base64_img = encode_image(str(image_path))
                    record = save_visual_record(doc_id, source_file, page_counter, image_path, visual_description)
                    vector_doc = create_vector_document(doc_id, source_file, page_counter, str(image_path), visual_description, base64_img)
                    
                    visual_records.append(record)
                    vector_documents.append(vector_doc)
                except Exception as e:
                    log.warning(f"Failed to extract image from DOCX: {e}")
        
        log.info(f"DOCX visual extraction: extracted {len(vector_documents)} visuals from {docx_path}")
    except ImportError:
        log.warning("python-docx not installed, skipping DOCX visual extraction")
    except Exception as e:
        log.error(f"DOCX visual extraction failed for {docx_path}: {e}")
    
    return visual_records, vector_documents


def process_pptx_visuals(pptx_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Extract and render slides from PPTX files as images and analyze them.
    """
    visual_records = []
    vector_documents = []
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation(pptx_path)
        source_file = os.path.basename(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Convert slide to image using python-pptx
                # First, save slide as temporary PDF using LibreOffice if available
                # Or use PIL to render the slide
                
                # Alternative approach: extract shapes and text
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                # If slide has images, extract them
                image_found = False
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture shape
                        try:
                            image = shape.image
                            image_filename = f"{doc_id}_slide_{slide_num}.png"
                            image_path = STORAGE_DIR / image_filename
                            
                            with open(image_path, "wb") as f:
                                f.write(image.blob)
                            
                            visual_description = analyze_image_with_gpt4o(str(image_path))
                            
                            if "NO_VISUAL_CONTENT" in visual_description:
                                if image_path.exists():
                                    os.remove(image_path)
                                continue
                            
                            image_found = True
                            base64_img = encode_image(str(image_path))
                            record = save_visual_record(doc_id, source_file, slide_num, image_path, visual_description)
                            vector_doc = create_vector_document(doc_id, source_file, slide_num, str(image_path), visual_description, base64_img)
                            
                            visual_records.append(record)
                            vector_documents.append(vector_doc)
                        except Exception as e:
                            log.warning(f"Failed to extract image from PPTX slide: {e}")
                
                if not image_found and slide_text:
                    # Try rendering slide to image if has content
                    try:
                        import subprocess
                        temp_pptx = STORAGE_DIR / f"temp_{doc_id}_slide_{slide_num}.pptx"
                        temp_pptx.parent.mkdir(parents=True, exist_ok=True)
                        
                        # Render using LibreOffice if available
                        result = subprocess.run(
                            [
                                "libreoffice", "--headless", "--convert-to", "pdf",
                                "--outdir", str(STORAGE_DIR), str(pptx_path)
                            ],
                            capture_output=True,
                            timeout=30
                        )
                        # This is advanced; for now just log
                        log.debug(f"PPTX slide {slide_num} rendering attempted")
                    except Exception as e:
                        log.debug(f"Could not render PPTX slide: {e}")
            except Exception as e:
                log.warning(f"Failed to process PPTX slide {slide_num}: {e}")
        
        log.info(f"PPTX visual extraction: extracted {len(vector_documents)} visuals from {pptx_path}")
    except ImportError:
        log.warning("python-pptx not installed, skipping PPTX visual extraction")
    except Exception as e:
        log.error(f"PPTX visual extraction failed for {pptx_path}: {e}")
    
    return visual_records, vector_documents


def process_xlsx_visuals(xlsx_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Extract charts and formatted data from XLSX files.
    """
    visual_records = []
    vector_documents = []
    
    try:
        from openpyxl import load_workbook
        from openpyxl.drawing.image import Image as XlsxImage
        
        workbook = load_workbook(xlsx_path)
        source_file = os.path.basename(xlsx_path)
        page_counter = 0
        
        for sheet_name in workbook.sheetnames:
            try:
                sheet = workbook[sheet_name]
                
                # Extract images/charts from sheet
                for drawing in sheet._drawing_list:
                    try:
                        if hasattr(drawing, 'image'):
                            page_counter += 1
                            image_data = drawing.image.blob
                            image_filename = f"{doc_id}_sheet_{sheet_name}_img_{page_counter}.png"
                            image_path = STORAGE_DIR / image_filename
                            
                            with open(image_path, "wb") as f:
                                f.write(image_data)
                            
                            visual_description = analyze_image_with_gpt4o(str(image_path))
                            
                            if "NO_VISUAL_CONTENT" in visual_description:
                                if image_path.exists():
                                    os.remove(image_path)
                                continue
                            
                            base64_img = encode_image(str(image_path))
                            record = save_visual_record(doc_id, source_file, page_counter, image_path, visual_description)
                            vector_doc = create_vector_document(doc_id, source_file, page_counter, str(image_path), visual_description, base64_img)
                            
                            visual_records.append(record)
                            vector_documents.append(vector_doc)
                    except Exception as e:
                        log.warning(f"Failed to extract chart from XLSX: {e}")
                
                # Optionally: render sheet as image for complex layouts
                try:
                    import subprocess
                    temp_pdf = STORAGE_DIR / f"temp_{doc_id}_{sheet_name}.pdf"
                    result = subprocess.run(
                        [
                            "libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", str(STORAGE_DIR), str(xlsx_path)
                        ],
                        capture_output=True,
                        timeout=30
                    )
                    log.debug(f"XLSX sheet {sheet_name} rendering attempted")
                except Exception as e:
                    log.debug(f"Could not render XLSX sheet: {e}")
            except Exception as e:
                log.warning(f"Failed to process XLSX sheet '{sheet_name}': {e}")
        
        log.info(f"XLSX visual extraction: extracted {len(vector_documents)} visuals from {xlsx_path}")
    except ImportError:
        log.warning("openpyxl not installed, skipping XLSX visual extraction")
    except Exception as e:
        log.error(f"XLSX visual extraction failed for {xlsx_path}: {e}")
    
    return visual_records, vector_documents


def process_image_visuals(image_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Process direct image files (PNG, JPG, JPEG) as visual evidence.
    """
    visual_records = []
    vector_documents = []

    try:
        source_file = os.path.basename(image_path)
        page_counter = 1

        # Copy/reference the image
        image_filename = f"{doc_id}_image_1.png"
        dest_path = STORAGE_DIR / image_filename

        try:
            from PIL import Image
            img = Image.open(image_path)
            img.save(str(dest_path), "PNG")
        except Exception:
            # Fallback: just copy the file
            import shutil
            shutil.copy(image_path, dest_path)

        # Analyze with GPT-4o
        visual_description = analyze_image_with_gpt4o(str(dest_path))

        if "NO_VISUAL_CONTENT" in visual_description:
            if dest_path.exists():
                os.remove(dest_path)
            log.info(f"Image file {source_file} contained no meaningful visual content")
            return [], []

        base64_img = encode_image(str(dest_path))
        record = save_visual_record(doc_id, source_file, page_counter, dest_path, visual_description)
        vector_doc = create_vector_document(doc_id, source_file, page_counter, str(dest_path), visual_description, base64_img)

        visual_records.append(record)
        vector_documents.append(vector_doc)

        log.info(f"Image visual extraction: extracted 1 visual from {image_path}")
    except Exception as e:
        log.error(f"Image visual extraction failed for {image_path}: {e}")

    return visual_records, vector_documents


def process_eml_visuals(eml_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Extract embedded images from EML (email) files.
    """
    visual_records = []
    vector_documents = []

    try:
        import email
        from email import policy
        from email.mime.multipart import MIMEMultipart

        source_file = os.path.basename(eml_path)
        page_counter = 0

        with open(eml_path, "rb") as f:
            msg = email.message_from_bytes(f.read(), policy=policy.default)

        # Extract images from multipart message
        for part in msg.iter_parts():
            if part.get_content_maintype() == "image":
                try:
                    page_counter += 1
                    image_data = part.get_payload(decode=True)
                    image_filename = f"{doc_id}_email_img_{page_counter}.png"
                    image_path = STORAGE_DIR / image_filename

                    with open(image_path, "wb") as f:
                        f.write(image_data)

                    visual_description = analyze_image_with_gpt4o(str(image_path))

                    if "NO_VISUAL_CONTENT" in visual_description:
                        if image_path.exists():
                            os.remove(image_path)
                        continue

                    base64_img = encode_image(str(image_path))
                    record = save_visual_record(doc_id, source_file, page_counter, image_path, visual_description)
                    vector_doc = create_vector_document(doc_id, source_file, page_counter, str(image_path), visual_description, base64_img)

                    visual_records.append(record)
                    vector_documents.append(vector_doc)
                except Exception as e:
                    log.warning(f"Failed to extract image from EML: {e}")

        log.info(f"EML visual extraction: extracted {len(vector_documents)} visuals from {eml_path}")
    except Exception as e:
        log.error(f"EML visual extraction failed for {eml_path}: {e}")

    return visual_records, vector_documents


def process_document_visuals(file_path: str, doc_id: str) -> Tuple[List[dict], List[Document]]:
    """
    Route visual extraction based on file format.
    Supports: PDF, DOCX, PPTX, XLSX, PNG, JPG, JPEG, EML
    """
    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == ".pdf":
        return process_pdf_visuals(file_path, doc_id)
    elif file_ext in (".docx", ".doc"):
        return process_docx_visuals(file_path, doc_id)
    elif file_ext in (".pptx", ".ppt"):
        return process_pptx_visuals(file_path, doc_id)
    elif file_ext in (".xlsx", ".xls"):
        return process_xlsx_visuals(file_path, doc_id)
    elif file_ext in (".png", ".jpg", ".jpeg"):
        return process_image_visuals(file_path, doc_id)
    elif file_ext == ".eml":
        return process_eml_visuals(file_path, doc_id)
    else:
        log.warning(f"Unsupported file format for visual extraction: {file_ext}")
        return [], []